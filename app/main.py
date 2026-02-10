import json
import os
import uuid
from datetime import datetime
from typing import List, Optional, Tuple

import httpx
from convertdate import hebrew
from dateutil import parser
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(BASE_DIR, "data")
GENERATED_DIR = os.path.join(os.path.dirname(BASE_DIR), "generated")
STATIC_DIR = os.path.join(os.path.dirname(BASE_DIR), "static")

with open(os.path.join(DATA_DIR, "etcbc_samples.json"), "r", encoding="utf-8") as fh:
    ETCBC_SAMPLES = json.load(fh)

with open(os.path.join(DATA_DIR, "congregation_calendar.json"), "r", encoding="utf-8") as fh:
    CONG_CALENDAR = json.load(fh)

FESTIVALS = [
    (1, 15, "Passover (Pesach)", "Celebrates redemption from Egypt and anticipates ultimate deliverance."),
    (1, 21, "Feast of Unleavened Bread", "Calls to remove leaven—symbolizing holiness and readiness."),
    (3, 6, "Shavuot (Pentecost)", "Remembers Torah giving and the Spirit's empowering."),
    (7, 1, "Rosh Hashanah", "Invites reflection, repentance, and attentiveness to God's voice."),
    (7, 10, "Yom Kippur", "Centers on atonement and God's mercy."),
    (7, 15, "Sukkot", "Highlights God's provision in wilderness journeys."),
    (9, 25, "Hanukkah", "Celebrates dedication and faithful witness."),
    (12, 14, "Purim", "Recounts God's hidden deliverance in Esther's story."),
]

COLOR_PALETTE = [
    (40, 75, 99),
    (142, 68, 173),
    (22, 160, 133),
    (192, 57, 43),
]


def ensure_generated_dir() -> None:
    os.makedirs(GENERATED_DIR, exist_ok=True)


def search_etcbc(query: str, limit: int = 5) -> List[dict]:
    """Attempt to search the ETCBC GitHub repository. Fallback to local samples."""
    headers = {"Accept": "application/vnd.github.v3+json"}
    url = "https://api.github.com/search/code"
    params = {"q": f"{query} repo:ETCBC/bhsa", "per_page": str(limit)}
    try:
        with httpx.Client(timeout=5.0) as client:
            response = client.get(url, params=params, headers=headers)
            if response.status_code == 200:
                data = response.json()
                items = []
                for item in data.get("items", [])[:limit]:
                    items.append(
                        {
                            "name": item.get("name"),
                            "path": item.get("path"),
                            "html_url": item.get("html_url"),
                            "repository": item.get("repository", {}).get("full_name", "ETCBC/bhsa"),
                        }
                    )
                if items:
                    return items
    except httpx.HTTPError:
        pass
    # fallback to local sample metadata
    fallback = []
    for sample in ETCBC_SAMPLES:
        fallback.append(
            {
                "name": sample["reference"],
                "path": f"samples/{sample['book'].lower()}.json",
                "html_url": "https://github.com/ETCBC/bhsa",
                "repository": "ETCBC/bhsa",
            }
        )
    return fallback[:limit]


def select_sample(topic: Optional[str], passage: Optional[str]) -> dict:
    normalized_topic = (topic or "").lower()
    normalized_passage = (passage or "").lower()
    for sample in ETCBC_SAMPLES:
        if normalized_passage and normalized_passage in sample["reference"].lower():
            return sample
    if normalized_topic:
        for sample in ETCBC_SAMPLES:
            if any(normalized_topic in theme.lower() for theme in sample["themes"]):
                return sample
    return ETCBC_SAMPLES[0]


def congregation_context(congregation_id: str, target_date: datetime) -> dict:
    record = CONG_CALENDAR.get(congregation_id) or CONG_CALENDAR.get("default", {})
    events = []
    for item in record.get("significant_dates", []):
        try:
            event_date = parser.isoparse(item["date"]).date()
        except (ValueError, TypeError):
            continue
        delta = abs((target_date.date() - event_date).days)
        if delta <= 21:
            events.append({"description": item["description"], "emphasis": item.get("emphasis"), "days_apart": delta})
    return {
        "name": record.get("name"),
        "location": record.get("location"),
        "values": record.get("values", []),
        "nearby_events": sorted(events, key=lambda e: e["days_apart"]),
    }


def hebrew_festival_matches(gregorian_date: datetime) -> List[dict]:
    hy, _, _ = hebrew.from_gregorian(gregorian_date.year, gregorian_date.month, gregorian_date.day)
    matches = []
    for offset in (-1, 0, 1):
        year = hy + offset
        for month, day, name, emphasis in FESTIVALS:
            try:
                g_year, g_month, g_day = hebrew.to_gregorian(year, month, day)
                festival_date = datetime(g_year, g_month, g_day)
            except ValueError:
                continue
            delta = abs((gregorian_date - festival_date).days)
            if delta <= 21:
                matches.append(
                    {
                        "festival": name,
                        "emphasis": emphasis,
                        "festival_date": festival_date.date().isoformat(),
                        "days_apart": delta,
                    }
                )
    matches.sort(key=lambda m: m["days_apart"])
    return matches


def build_introduction(sample: dict, audience: Optional[str], occasion: Optional[str], festival_info: List[dict], congregation_info: dict) -> str:
    hook = f"Imagine standing where {sample['reference']} first unfolded—hearing the Hebrew cadence of {sample['hebrew_focus']} inviting trust." 
    audience_line = f"For {audience}, " if audience else ""
    festival_line = ""
    if festival_info:
        nearest = festival_info[0]
        festival_line = (
            f" We gather with {nearest['festival']} approaching ({nearest['festival_date']}), a season inviting {nearest['emphasis'].lower()}."
        )
    congregation_line = ""
    if congregation_info.get("nearby_events"):
        first = congregation_info["nearby_events"][0]
        congregation_line = (
            f" Our own community prepares for {first['description']} in {first['days_apart']} days, aligning hearts toward {first['emphasis'].lower()}."
        )
    occasion_line = f" In this {occasion.lower()} we are called to listen afresh." if occasion else ""
    return (
        f"{hook}{festival_line}{congregation_line} {audience_line}God's word speaks with precision and promise.{occasion_line}"
    ).strip()


def build_conclusion(sample: dict, introduction: str) -> str:
    focus = sample["hebrew_focus"]
    return (
        f"The same cadence that opened our time—{focus}—now sends us. Let the insights we traced move from study to practice: "+
        "embrace God's invitation, embody covenantal blessing, and walk toward Christlike transformation together."
    )


def construct_sections(sample: dict, estimated_minutes: int) -> List[dict]:
    base_points = [
        {
            "title": "Textual Horizon",
            "content": (
                f"{sample['reference']} anchors the lesson. Key Hebrew focus: {sample['hebrew_focus']} ({sample['translation']}). "
                f"Morphology: {sample['morphology']['part_of_speech']} rooted in {sample['morphology']['root']}."
            ),
            "exegetical": sample["exegetical_notes"],
            "application": "Trace the flow of the passage, inviting listeners to inhabit the narrative movement.",
        },
        {
            "title": "Linguistic Insights",
            "content": "Lexical themes emerge, amplifying covenantal movement.",
            "exegetical": sample["lexical_insights"],
            "application": "Highlight how the Hebrew terms reshape imagination and discipleship practices.",
        },
        {
            "title": "Formation Pathways",
            "content": "Move from exegesis to embodied action.",
            "exegetical": [
                "Map the passage's structure to contemporary rhythms (gathering, scattering, serving).",
                "Invite testimonies or reflective prayer that echo the text's movement.",
            ],
            "application": "Provide concrete steps for the congregation to live the text this week.",
        },
    ]
    if estimated_minutes < 25:
        # condense by merging second and third points
        condensed = base_points[:1]
        condensed.append(
            {
                "title": "Concise Insight",
                "content": (
                    f"In limited time, emphasize the pivot: {sample['hebrew_focus']} propels us toward faithful obedience."
                ),
                "exegetical": sample["lexical_insights"][:1],
                "application": "Offer one spiritual practice and one communal action.",
            }
        )
        return condensed
    return base_points


def estimate_runtime(estimated_minutes: int, interpreted: bool) -> int:
    if interpreted:
        return max(10, int(estimated_minutes * 0.65))
    return estimated_minutes


def make_slide_palette(index: int) -> Tuple[int, int, int]:
    return COLOR_PALETTE[index % len(COLOR_PALETTE)]


def create_slides(introduction: str, sections: List[dict], conclusion: str) -> List[dict]:
    slides = [
        {
            "title": "Opening Story",
            "bullets": [introduction],
            "notes": "Welcome, frame the occasion, and read the passage aloud.",
        }
    ]
    for idx, section in enumerate(sections, start=1):
        slides.append(
            {
                "title": f"{idx}. {section['title']}",
                "bullets": [section["content"], *section["exegetical"], section["application"]],
                "notes": "Guide discussion; invite observations and response.",
            }
        )
    slides.append(
        {
            "title": "Sending Charge",
            "bullets": [conclusion],
            "notes": "Summarize commitments and pray a commissioning blessing.",
        }
    )
    return slides


def create_pptx(slides: List[dict]) -> str:
    ensure_generated_dir()
    presentation = Presentation()
    title_layout = presentation.slide_layouts[1]
    for idx, slide_data in enumerate(slides):
        slide = presentation.slides.add_slide(title_layout)
        background_color = make_slide_palette(idx)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*background_color)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = slide_data["title"]
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        body_shape = shapes.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        for idx, bullet in enumerate(slide_data["bullets"]):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.level = 0 if idx == 0 else 1
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data["notes"]
    file_id = uuid.uuid4().hex
    file_path = os.path.join(GENERATED_DIR, f"lesson_{file_id}.pptx")
    presentation.save(file_path)
    return file_path


class LessonRequest(BaseModel):
    audience: Optional[str]
    occasion: Optional[str]
    date: Optional[str]
    topic: Optional[str]
    passage: Optional[str]
    lesson_type: str = Field(regex="^(expository|topical|bible_study|personal)$")
    estimated_minutes: int = Field(default=35, ge=10, le=120)
    interpreted: bool = False
    congregation_id: str = "default"


class LessonResponse(BaseModel):
    lesson: dict
    festivals: List[dict]
    congregation: dict
    github_sources: List[dict]
    runtime_minutes: int
    pptx_download: str


app = FastAPI(title="ETCBC Lesson Designer")
app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")


@app.get("/")
async def root():
    return FileResponse(os.path.join(STATIC_DIR, "index.html"))


@app.post("/api/generate", response_model=LessonResponse)
async def generate_lesson(payload: LessonRequest):
    if not payload.topic and not payload.passage:
        raise HTTPException(status_code=400, detail="Topic or passage required")

    target_date = datetime.utcnow()
    if payload.date:
        try:
            target_date = parser.isoparse(payload.date)
        except (ValueError, TypeError):
            raise HTTPException(status_code=400, detail="Invalid date format")

    sample = select_sample(payload.topic, payload.passage)

    festivals = hebrew_festival_matches(target_date)
    congregation = congregation_context(payload.congregation_id, target_date)
    introduction = build_introduction(sample, payload.audience, payload.occasion, festivals, congregation)
    conclusion = build_conclusion(sample, introduction)
    sections = construct_sections(sample, payload.estimated_minutes)
    slides = create_slides(introduction, sections, conclusion)
    pptx_path = create_pptx(slides)
    runtime_minutes = estimate_runtime(payload.estimated_minutes, payload.interpreted)

    github_query = payload.passage or payload.topic or sample["reference"]
    github_sources = search_etcbc(github_query)

    lesson = {
        "title": f"{sample['reference']} — {payload.lesson_type.replace('_', ' ').title()}",
        "introduction": introduction,
        "conclusion": conclusion,
        "sections": sections,
        "slides": slides,
        "canvas": {
            "introduction": introduction,
            "sections": sections,
            "conclusion": conclusion,
        },
        "hebrew_focus": sample["hebrew_focus"],
        "morphology": sample["morphology"],
        "themes": sample["themes"],
    }

    return LessonResponse(
        lesson=lesson,
        festivals=festivals,
        congregation=congregation,
        github_sources=github_sources,
        runtime_minutes=runtime_minutes,
        pptx_download=f"/api/pptx/{os.path.basename(pptx_path)}",
    )


@app.get("/api/pptx/{filename}")
async def download_pptx(filename: str):
    file_path = os.path.join(GENERATED_DIR, filename)
    if not os.path.isfile(file_path):
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(file_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)
